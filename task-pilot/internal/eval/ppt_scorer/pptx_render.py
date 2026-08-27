#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""纯 Python pptx → PNG 渲染器,给 render_deck.py 调用。

设计目标:只依赖 pip 包(python-pptx + Pillow),不 fork 任何系统二进制
(soffice/libreoffice/pdftoppm 等一概不需要),因此评测执行器镜像只要能
`pip install python-pptx Pillow` 就能跑,不再受镜像是否预装 LibreOffice 摆布。

覆盖范围(够 VLM 判排版缺陷即可,不追求逐像素还原 PowerPoint):
  - 幻灯片背景色(纯色;渐变/图片背景取首个渐变色或跳过画纯灰兜底)
  - 自由形状/矩形/椭圆等 autoshape:纯色填充 + 描边
  - 连接线(直线):按起止点画线
  - 文本框/占位符文本:逐 paragraph/run 绘制,尊重字号、粗体、斜体、颜色、
    水平对齐;字体统一用内置 CJK 兜底字体(不依赖系统装了什么字体)
  - 表格:画网格线 + 逐格文字(自动换行、超出画省略号提示,便于 VLM 识别
    table_bad/overflow)
  - 图片(p:pic):直接解出 blob 贴到对应矩形区域(保持宽高比居中裁剪)
  - 图表(chart)/其它未识别对象:画一个带类型标签的占位框,不尝试还原具体
    图表样式(VLM 评审只需要知道"这里有个图表"这一事实)
  - 组合形状(group):递归处理子形状,按 chOff/chExt → off/ext 做子坐标系
    换算

已知不还原的细节(不影响排版缺陷判断的部分):渐变/图片填充只取近似色、
阴影/发光/3D 效果、动画、精确字体度量(用近似字宽估算换行,可能与真实
PowerPoint 换行位置有 1-2 字符出入)。
"""
from __future__ import annotations

import io
import os
import urllib.request
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont

EMU_PER_INCH = 914400
DEFAULT_DPI = 110

# 内置渲染字体:优先复用同目录下随 PPTScorerFilenames 一起下发的 CJK 字体文件
# (下发机制是扁平文件列表,统一落在 tmp/ppt_scorer/ 里,不支持子目录,故用同目录
# 平铺文件名而非 fonts/ 子目录),再退化到常见镜像可能预装的系统 CJK 字体路径,
# 最后才落到 Pillow 默认字体(不含中文字形,中文会渲染成方框/缺字,仅兜底英文场景)。
_HERE = Path(__file__).resolve().parent
_FONT_CANDIDATES = [
    _HERE / "NotoSansSC-Regular.otf",
    _HERE / "NotoSansSC-Regular.ttf",
    Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
    Path("/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc"),
    Path("/usr/share/fonts/truetype/wqy/wqy-microhei.ttc"),
    Path("/usr/share/fonts/truetype/arphic/uming.ttc"),
]
_FONT_ENV = os.environ.get("PPT_RENDER_FONT")
if _FONT_ENV:
    _FONT_CANDIDATES.insert(0, Path(_FONT_ENV))

_FONT_CACHE: dict[int, ImageFont.FreeTypeFont] = {}
_FONT_PATH: str | None = None


def _resolve_font_path() -> str | None:
    global _FONT_PATH
    if _FONT_PATH is not None:
        return _FONT_PATH or None
    for cand in _FONT_CANDIDATES:
        if cand.is_file():
            _FONT_PATH = str(cand)
            return _FONT_PATH
    _FONT_PATH = ""
    return None


def _font(size_px: int) -> ImageFont.ImageFont:
    size_px = max(int(size_px), 6)
    cached = _FONT_CACHE.get(size_px)
    if cached is not None:
        return cached
    path = _resolve_font_path()
    try:
        if path:
            f = ImageFont.truetype(path, size_px)
        else:
            f = ImageFont.load_default(size=size_px)
    except Exception:
        f = ImageFont.load_default()
    _FONT_CACHE[size_px] = f
    return f


def _emu_to_px(emu: int | None, dpi: int) -> int:
    if emu is None:
        return 0
    return round(emu / EMU_PER_INCH * dpi)


def _rgb_from_color(color, theme_colors: dict[str, tuple[int, int, int]] | None) -> tuple[int, int, int] | None:
    """尽力从 python-pptx ColorFormat 取出 RGB;拿不到(主题色引用等)返回 None。"""
    if color is None:
        return None
    try:
        ctype = color.type
    except Exception:
        return None
    if ctype is None:
        return None
    try:
        from pptx.enum.dml import MSO_COLOR_TYPE
    except Exception:
        MSO_COLOR_TYPE = None  # noqa: N806
    if MSO_COLOR_TYPE is not None and ctype == MSO_COLOR_TYPE.RGB:
        rgb = color.rgb
        return (rgb[0], rgb[1], rgb[2]) if rgb is not None else None
    if MSO_COLOR_TYPE is not None and ctype == MSO_COLOR_TYPE.SCHEME and theme_colors:
        key = str(color.theme_color).split(".")[-1]
        return theme_colors.get(key)
    return None


def _extract_theme_colors(prs) -> dict[str, tuple[int, int, int]]:
    """从第一个 slide master 的 theme 里取配色方案,供 scheme color 兜底。"""
    out: dict[str, tuple[int, int, int]] = {}
    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        from pptx.oxml import parse_xml
        el = parse_xml(theme_part.blob)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        scheme = el.find(".//a:clrScheme", ns)
        if scheme is None:
            return out
        name_map = {
            "dk1": "DARK_1", "lt1": "LIGHT_1", "dk2": "DARK_2", "lt2": "LIGHT_2",
            "accent1": "ACCENT_1", "accent2": "ACCENT_2", "accent3": "ACCENT_3",
            "accent4": "ACCENT_4", "accent5": "ACCENT_5", "accent6": "ACCENT_6",
            "hlink": "HYPERLINK", "folHlink": "FOLLOWED_HYPERLINK",
        }
        for child in scheme:
            tag = child.tag.split("}")[-1]
            key = name_map.get(tag)
            if not key:
                continue
            srgb = child.find("a:srgbClr", ns)
            if srgb is not None and srgb.get("val"):
                v = srgb.get("val")
                out[key] = (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
            else:
                sysclr = child.find("a:sysClr", ns)
                if sysclr is not None and sysclr.get("lastClr"):
                    v = sysclr.get("lastClr")
                    out[key] = (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
    except Exception:
        pass
    return out


def _shape_fill_rgb(shape, theme_colors) -> tuple[int, int, int] | None:
    try:
        fill = shape.fill
    except Exception:
        return None
    try:
        ftype = fill.type
    except Exception:
        return None
    if ftype is None:
        return None
    from pptx.enum.dml import MSO_FILL_TYPE
    if ftype == MSO_FILL_TYPE.SOLID:
        return _rgb_from_color(fill.fore_color, theme_colors)
    if ftype == MSO_FILL_TYPE.GRADIENT:
        try:
            stops = list(fill.gradient_stops)
            if stops:
                return _rgb_from_color(stops[0].color, theme_colors)
        except Exception:
            return None
    return None


def _shape_line_rgb(shape, theme_colors) -> tuple[tuple[int, int, int] | None, int]:
    try:
        line = shape.line
        rgb = _rgb_from_color(line.color, theme_colors)
        width_px = max(1, round((line.width or 0) / 12700))  # EMU->pt(12700/pt)->px近似(1pt≈1px@96dpi)
        return rgb, width_px
    except Exception:
        return None, 1


def _draw_rect_or_ellipse(draw: ImageDraw.ImageDraw, shape, box, theme_colors) -> None:
    x0, y0, x1, y1 = box
    if x1 <= x0 or y1 <= y0:
        return
    fill = _shape_fill_rgb(shape, theme_colors)
    line_rgb, line_w = _shape_line_rgb(shape, theme_colors)
    fill_rgba = (*fill, 255) if fill else None
    outline_rgba = (*line_rgb, 255) if line_rgb else None

    prst = None
    try:
        prst = shape.auto_shape_type
    except Exception:
        prst = None
    is_ellipse = False
    try:
        from pptx.enum.shapes import MSO_SHAPE
        is_ellipse = prst in (MSO_SHAPE.OVAL,)
    except Exception:
        pass

    if is_ellipse:
        draw.ellipse(box, fill=fill_rgba, outline=outline_rgba, width=line_w if outline_rgba else 0)
    else:
        draw.rectangle(box, fill=fill_rgba, outline=outline_rgba, width=line_w if outline_rgba else 0)


def _draw_line_shape(draw: ImageDraw.ImageDraw, shape, box, theme_colors) -> None:
    x0, y0, x1, y1 = box
    line_rgb, line_w = _shape_line_rgb(shape, theme_colors)
    color = (*line_rgb, 255) if line_rgb else (120, 120, 120, 255)
    draw.line([(x0, y0), (x1, y1)], fill=color, width=max(1, line_w))


def _wrap_text(text: str, font: ImageFont.ImageFont, max_width: int, draw: ImageDraw.ImageDraw) -> list[str]:
    if max_width <= 0:
        return [text]
    lines: list[str] = []
    cur = ""
    for ch in text:
        trial = cur + ch
        w = draw.textlength(trial, font=font)
        if w > max_width and cur:
            lines.append(cur)
            cur = ch
        else:
            cur = trial
    if cur:
        lines.append(cur)
    return lines or [""]


def _pt_to_px(pt_emu_or_pt, dpi: int) -> int:
    """python-pptx font.size 是 Emu 的 Pt 子类,直接 .pt 拿磅数;换算到像素按 dpi/72。"""
    try:
        pt = pt_emu_or_pt.pt
    except Exception:
        pt = 18
    return max(6, round(pt * dpi / 72))


def _draw_text_frame(draw: ImageDraw.ImageDraw, shape, box, theme_colors, dpi: int) -> None:
    try:
        tf = shape.text_frame
    except Exception:
        return
    x0, y0, x1, y1 = box
    pad = max(2, round(0.05 * dpi))
    inner_x0, inner_y0, inner_x1 = x0 + pad, y0 + pad, x1 - pad
    max_width = max(1, inner_x1 - inner_x0)

    from pptx.enum.text import PP_ALIGN

    cur_y = inner_y0
    for para in tf.paragraphs:
        runs = list(para.runs) or [None]
        # 段落整体按第一个有内容的 run 定字号,估算行高;多 run 混排时逐 run 变宽度但共享行高基线。
        sizes = [r.font.size for r in para.runs if r is not None and r.font.size is not None]
        font_px = _pt_to_px(sizes[0], dpi) if sizes else round(0.18 * dpi)
        base_font = _font(font_px)
        line_h = round(font_px * 1.3)

        # 拼出整段文字用于换行估计,再按 run 顺序重新切片着色绘制(简化:每个 run 单独换行绘制,
        # run 之间不跨行拼接——对多数「一段一 run」或「run 边界=语义边界」的 deck 足够,复杂
        # 混排场景可能行尾留白,不影响排版缺陷类判断)。
        align = para.alignment
        for run in (para.runs or [_EmptyRun()]):
            text = getattr(run, "text", "") or ""
            if text == "":
                continue
            size_px = _pt_to_px(run.font.size, dpi) if run.font.size else font_px
            font = _font(size_px)
            color = _rgb_from_color(run.font.color, theme_colors) if run.font.color else None
            fill = (*color, 255) if color else (30, 30, 30, 255)
            lines = _wrap_text(text, font, max_width, draw)
            for ln in lines:
                if cur_y > y1:
                    break
                w = draw.textlength(ln, font=font)
                if align == PP_ALIGN.CENTER:
                    lx = inner_x0 + max(0, (max_width - w) / 2)
                elif align == PP_ALIGN.RIGHT:
                    lx = inner_x0 + max(0, max_width - w)
                else:
                    lx = inner_x0
                try:
                    if getattr(run.font, "bold", False):
                        draw.text((lx, cur_y), ln, font=font, fill=fill, stroke_width=1, stroke_fill=fill)
                    else:
                        draw.text((lx, cur_y), ln, font=font, fill=fill)
                except Exception:
                    draw.text((lx, cur_y), ln, font=font, fill=fill)
                cur_y += line_h
        cur_y += round(line_h * 0.15)  # 段间距


class _EmptyRun:
    text = ""
    font = None


def _draw_table(draw: ImageDraw.ImageDraw, shape, box, theme_colors, dpi: int) -> None:
    try:
        table = shape.table
    except Exception:
        return
    x0, y0, x1, y1 = box
    n_rows = len(table.rows)
    n_cols = len(table.columns)
    if n_rows == 0 or n_cols == 0:
        return

    total_w = x1 - x0
    total_h = y1 - y0
    col_ws = [c.width for c in table.columns]
    tw = sum(col_ws) or 1
    row_hs = [r.height for r in table.rows]
    th = sum(row_hs) or 1

    xs = [x0]
    for w in col_ws:
        xs.append(xs[-1] + total_w * (w / tw))
    ys = [y0]
    for h in row_hs:
        ys.append(ys[-1] + total_h * (h / th))

    grid_color = (120, 120, 120, 255)
    for xv in xs:
        draw.line([(xv, y0), (xv, y1)], fill=grid_color, width=1)
    for yv in ys:
        draw.line([(x0, yv), (x1, yv)], fill=grid_color, width=1)

    font_px = max(8, round(0.12 * dpi))
    font = _font(font_px)
    pad = 3
    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cx0, cy0, cx1, cy1 = xs[c] + pad, ys[r] + pad, xs[c + 1] - pad, ys[r + 1] - pad
            if cx1 <= cx0 or cy1 <= cy0:
                continue
            text = (cell.text or "").strip()
            if not text:
                continue
            lines = _wrap_text(text, font, max(1, cx1 - cx0), draw)
            cy = cy0
            for ln in lines:
                if cy + font_px > cy1:
                    # 超出单元格高度:画省略号提示截断(table_bad/overflow 信号)
                    draw.text((cx0, max(cy0, cy1 - font_px)), "...", font=font, fill=(200, 40, 40, 255))
                    break
                draw.text((cx0, cy), ln, font=font, fill=(20, 20, 20, 255))
                cy += round(font_px * 1.25)


def _draw_picture(draw: ImageDraw.ImageDraw, canvas: Image.Image, shape, box) -> None:
    x0, y0, x1, y1 = box
    w, h = max(1, x1 - x0), max(1, y1 - y0)
    try:
        blob = shape.image.blob
        img = Image.open(io.BytesIO(blob)).convert("RGBA")
    except Exception:
        draw.rectangle(box, fill=(230, 230, 230, 255), outline=(150, 150, 150, 255))
        draw.text((x0 + 4, y0 + 4), "[image]", font=_font(12), fill=(120, 120, 120, 255))
        return
    # 保持宽高比,居中裁剪填满目标框(近似 PowerPoint 默认的裁剪表现)。
    src_ratio = img.width / img.height
    dst_ratio = w / h
    if src_ratio > dst_ratio:
        new_h = img.height
        new_w = round(new_h * dst_ratio)
        left = round((img.width - new_w) / 2)
        img = img.crop((left, 0, left + new_w, new_h))
    else:
        new_w = img.width
        new_h = round(new_w / dst_ratio)
        top = round((img.height - new_h) / 2)
        img = img.crop((0, top, new_w, top + new_h))
    img = img.resize((w, h))
    canvas.alpha_composite(img, (x0, y0))


def _draw_placeholder_box(draw: ImageDraw.ImageDraw, box, label: str) -> None:
    x0, y0, x1, y1 = box
    if x1 <= x0 or y1 <= y0:
        return
    draw.rectangle(box, fill=(245, 245, 245, 255), outline=(160, 160, 160, 255), width=1)
    font = _font(max(10, round((y1 - y0) * 0.08)))
    draw.text((x0 + 6, y0 + 6), label, font=font, fill=(140, 140, 140, 255))


def _child_to_parent_box(off, ext, ch_off, ch_ext, parent_box) -> tuple[int, int, int, int]:
    """把 group 内子形状的 chOff/chExt 坐标系换算到父 box 的像素坐标。"""
    px0, py0, px1, py1 = parent_box
    pw, ph = max(1, px1 - px0), max(1, py1 - py0)
    cw, ch = max(1, ch_ext[0]), max(1, ch_ext[1])
    sx = pw / cw
    sy = ph / ch
    x0 = px0 + (off[0] - ch_off[0]) * sx
    y0 = py0 + (off[1] - ch_off[1]) * sy
    x1 = x0 + ext[0] * sx
    y1 = y0 + ext[1] * sy
    return round(x0), round(y0), round(x1), round(y1)


def _group_ch_offset_extent(group_shape) -> tuple[tuple[int, int], tuple[int, int]] | None:
    try:
        xfrm = group_shape._element.grpSpPr.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm"
        )
        if xfrm is None:
            return None
        chOff = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}chOff")
        chExt = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}chExt")
        if chOff is None or chExt is None:
            return None
        return (int(chOff.get("x")), int(chOff.get("y"))), (int(chExt.get("cx")), int(chExt.get("cy")))
    except Exception:
        return None


def _shape_box(shape, dpi: int) -> tuple[int, int, int, int] | None:
    try:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
    except Exception:
        return None
    if left is None or top is None or width is None or height is None:
        return None
    x0 = _emu_to_px(left, dpi)
    y0 = _emu_to_px(top, dpi)
    x1 = _emu_to_px(left + width, dpi)
    y1 = _emu_to_px(top + height, dpi)
    return x0, y0, x1, y1


def _draw_shape(draw: ImageDraw.ImageDraw, canvas: Image.Image, shape, box, theme_colors, dpi: int) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        stype = shape.shape_type
    except Exception:
        stype = None

    if stype == MSO_SHAPE_TYPE.GROUP:
        _draw_group(draw, canvas, shape, box, theme_colors, dpi)
        return
    if stype == MSO_SHAPE_TYPE.PICTURE or stype == MSO_SHAPE_TYPE.LINKED_PICTURE:
        _draw_picture(draw, canvas, shape, box)
        return
    if getattr(shape, "has_table", False):
        _draw_table(draw, shape, box, theme_colors, dpi)
        return
    if stype == MSO_SHAPE_TYPE.CHART:
        _draw_placeholder_box(draw, box, "[chart]")
        return
    if stype == MSO_SHAPE_TYPE.LINE:
        _draw_line_shape(draw, shape, box, theme_colors)
        return
    if stype in (
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
        MSO_SHAPE_TYPE.MEDIA,
        MSO_SHAPE_TYPE.DIAGRAM,
        MSO_SHAPE_TYPE.IGX_GRAPHIC,
    ):
        _draw_placeholder_box(draw, box, f"[{stype}]")
        return

    # AUTO_SHAPE / TEXT_BOX / PLACEHOLDER 等:先画填充/描边形状,再叠文字。
    _draw_rect_or_ellipse(draw, shape, box, theme_colors)
    if getattr(shape, "has_text_frame", False):
        _draw_text_frame(draw, shape, box, theme_colors, dpi)


def _draw_group(draw, canvas, group_shape, parent_box, theme_colors, dpi) -> None:
    ch = _group_ch_offset_extent(group_shape)
    for sub in group_shape.shapes:
        sub_box = _shape_box(sub, dpi)
        if sub_box is None:
            continue
        if ch is not None:
            ch_off, ch_ext = ch
            box = _child_to_parent_box(
                (sub.left, sub.top), (sub.width, sub.height), ch_off, ch_ext, parent_box
            )
        else:
            box = sub_box
        _draw_shape(draw, canvas, sub, box, theme_colors, dpi)


def _slide_background_rgb(slide, theme_colors) -> tuple[int, int, int]:
    try:
        fill = slide.background.fill
        if fill.type is not None:
            from pptx.enum.dml import MSO_FILL_TYPE
            if fill.type == MSO_FILL_TYPE.SOLID:
                rgb = _rgb_from_color(fill.fore_color, theme_colors)
                if rgb:
                    return rgb
    except Exception:
        pass
    return (255, 255, 255)


def render_pptx_to_pngs(pptx_path: str, out_dir: str, dpi: int = DEFAULT_DPI) -> list[str]:
    from pptx import Presentation

    os.makedirs(out_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    theme_colors = _extract_theme_colors(prs)

    page_w = _emu_to_px(prs.slide_width, dpi)
    page_h = _emu_to_px(prs.slide_height, dpi)
    page_w = max(page_w, 10)
    page_h = max(page_h, 10)

    pages: list[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        bg = _slide_background_rgb(slide, theme_colors)
        canvas = Image.new("RGBA", (page_w, page_h), (*bg, 255))
        draw = ImageDraw.Draw(canvas)

        for shape in slide.shapes:
            box = _shape_box(shape, dpi)
            if box is None:
                continue
            try:
                _draw_shape(draw, canvas, shape, box, theme_colors, dpi)
            except Exception:
                # 单个形状渲染失败不应拖垮整页;跳过继续画其它形状。
                continue

        out_path = os.path.join(out_dir, f"p{idx:03d}.png")
        canvas.convert("RGB").save(out_path, "PNG")
        pages.append(out_path)

    return pages
